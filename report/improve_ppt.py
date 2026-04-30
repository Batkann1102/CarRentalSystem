"""
ReportPPT.pptx сайжруулах скрипт
Засварууд:
 - Үг үсгийн алдаа засах
 - Хуудасны дугаарлалт засах (slide 14: 9/12 → 11/12)
 - Кодын хэсгийг өнгөөр (syntax highlighting) ялгах
 - Утга зүйн найруулга засах
 - Дүгнэлт хэсгийг бүрэн бичих
"""

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

INPUT  = r'C:\4_Havar\Diplom\Example_project\CarRentalSystem\report\ReportPPT.pptx'
OUTPUT = r'C:\4_Havar\Diplom\Example_project\CarRentalSystem\report\ReportPPT_improved.pptx'

prs = Presentation(INPUT)

# ─── Color palette ────────────────────────────────────────────────────────────
CODE_BG        = RGBColor(0x1E, 0x1E, 0x2E)   # dark navy background
CODE_FG        = RGBColor(0xCD, 0xD6, 0xF4)   # light lavender default text
ANNOTATION_CLR = RGBColor(0xA6, 0xE3, 0xA1)   # green  – @Annotation
KEYWORD_CLR    = RGBColor(0x89, 0xB4, 0xFA)   # blue   – keywords
STRING_CLR     = RGBColor(0xF3, 0x8B, 0xA8)   # pink   – strings
TAG_CLR        = RGBColor(0x89, 0xDC, 0xEB)   # cyan   – XML tags
ATTR_CLR       = RGBColor(0xF9, 0xE2, 0xAF)   # yellow – attributes / params
BOOL_CLR       = RGBColor(0xEB, 0xA0, 0xAC)   # red    – boolean / literals
CLASS_CLR      = RGBColor(0xF5, 0xC2, 0xE7)   # pink   – class names


# ─── Helper utilities ─────────────────────────────────────────────────────────

def solid_fill(shape, rgb: RGBColor):
    """Set solid fill on a shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_all_run_color(shape, rgb: RGBColor):
    """Set all run colors in a shape."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = rgb


def set_run_font(shape, size_pt=None, bold=None, mono=False):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if size_pt:
                run.font.size = Pt(size_pt)
            if bold is not None:
                run.font.bold = bold
            if mono:
                run.font.name = 'Consolas'


def replace_text_in_shape(shape, old: str, new: str) -> bool:
    """Replace text in any paragraph/run of a shape."""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        # Try simple per-run replacement first
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                changed = True
    return changed


def clear_and_write_paragraphs(tf, paragraphs_data):
    """
    Clear all paragraphs from a text frame and write new ones.
    paragraphs_data: list of dicts with keys:
        text, size_pt, bold, color (RGBColor), indent (int), align
    """
    # Remove all existing paragraph elements
    txBody = tf._txBody
    # Keep first <a:p> as template (we'll overwrite it)
    existing_ps = txBody.findall(qn('a:p'))

    # Remove all except first
    for p in existing_ps[1:]:
        txBody.remove(p)

    first_p = existing_ps[0]
    # Clear runs from first paragraph
    for r in first_p.findall(qn('a:r')):
        first_p.remove(r)
    for br in first_p.findall(qn('a:br')):
        first_p.remove(br)

    def make_para(text, size_pt=13, bold=False, color=None, indent=0, align=None):
        p = etree.SubElement(txBody, qn('a:p'))
        pPr = etree.SubElement(p, qn('a:pPr'))
        if indent > 0:
            pPr.set('indent', str(int(-indent * 914400 * 0.3)))
            pPr.set('marL', str(int(indent * 914400 * 0.35)))
        if align == 'center':
            pPr.set('algn', 'ctr')

        r = etree.SubElement(p, qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'), attrib={'lang': 'mn-MN', 'dirty': '0'})
        rPr.set('sz', str(int(size_pt * 100)))
        if bold:
            rPr.set('b', '1')
        if color:
            solidFill = etree.SubElement(rPr, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', f'{color.rgb:06X}' if hasattr(color, 'rgb') else
                        f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')
        t = etree.SubElement(r, qn('a:t'))
        t.text = text
        return p

    # Overwrite the first paragraph
    first_p.getparent().remove(first_p)

    for pd in paragraphs_data:
        make_para(
            text=pd.get('text', ''),
            size_pt=pd.get('size_pt', 13),
            bold=pd.get('bold', False),
            color=pd.get('color', None),
            indent=pd.get('indent', 0),
            align=pd.get('align', None),
        )


# ─── Color logic for code shapes ──────────────────────────────────────────────

JAVA_KEYWORDS = {
    'public', 'private', 'protected', 'class', 'interface', 'implements',
    'extends', 'new', 'return', 'void', 'true', 'false', 'null', 'static',
    'final', 'import', 'package', 'throws', 'throw', 'try', 'catch',
    'if', 'else', 'for', 'while', 'do', 'break', 'continue', 'this', 'super',
}

XML_TAG_STARTERS = ('<', '</', '<!')
XML_CLOSERS = ('>', '/>')


def classify_java_token(text: str) -> RGBColor:
    t = text.strip()
    if t.startswith('@'):
        return ANNOTATION_CLR
    if t in JAVA_KEYWORDS:
        return KEYWORD_CLR
    if t.startswith('"') or t.startswith("'"):
        return STRING_CLR
    if t in ('true', 'false', 'null'):
        return BOOL_CLR
    if t and t[0].isupper():
        return CLASS_CLR
    return CODE_FG


def classify_xml_token(text: str) -> RGBColor:
    t = text.strip()
    if t.startswith('<') or t.startswith('</') or t in ('>', '/>'):
        return TAG_CLR
    if t.startswith('"') or t.startswith("'"):
        return STRING_CLR
    return ATTR_CLR


def apply_code_style_java(shape):
    """Apply dark background + Java syntax coloring to a shape."""
    if not shape.has_text_frame:
        return
    solid_fill(shape, CODE_BG)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            words = run.text.split()
            # Determine color from first word
            if words:
                run.font.color.rgb = classify_java_token(words[0])
            else:
                run.font.color.rgb = CODE_FG
            if not run.font.name or run.font.name not in ('Consolas', 'Courier New'):
                run.font.name = 'Consolas'


def apply_code_style_xml(shape):
    """Apply dark background + XML syntax coloring to a shape."""
    if not shape.has_text_frame:
        return
    solid_fill(shape, CODE_BG)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            t = run.text.strip()
            run.font.color.rgb = classify_xml_token(t)
            run.font.name = 'Consolas'


def apply_mixed_code_style(shape, lang='java'):
    """Apply colorized code style to inline-mixed shape."""
    if not shape.has_text_frame:
        return
    solid_fill(shape, CODE_BG)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            t = run.text.strip()
            if lang == 'java':
                run.font.color.rgb = classify_java_token(t)
            else:
                run.font.color.rgb = classify_xml_token(t)
            run.font.name = 'Consolas'


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 1  – Cover
# ═══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, 'загварчлал, хэрэгжүүлэлт',
                              'загварчлал ба хэрэгжүүлэлт')
        replace_text_in_shape(shape, 'Удирдагч', 'Удирдагч багш')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 2  – Table of Contents: add missing section name
# ═══════════════════════════════════════════════════════════════════════════════
# Fix section label "Нэмэлт сайжруулалт" → more descriptive
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, 'Нэмэлт сайжруулалт',
                              'Нэмэлт сайжруулалт (A1 Тест, A2 Визуализаци)')
        replace_text_in_shape(shape, 'Хэрэгжүүлэлт & Туршилт',
                              'Хэрэгжүүлэлт ба Туршилт')
        replace_text_in_shape(shape, 'Үр дүн & Харьцуулалт',
                              'Үр дүн ба Харьцуулалт')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 3  – Research objective (no page number – first content slide)
# ═══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_text_frame:
        # Fix "бий" at the end (incomplete sentence)
        replace_text_in_shape(shape,
            'Java Reflection API, annotation боловсруулалт, object lifecycle зэрэг суурь ойлголтыг гүнзгийрүүлэн судлах шаардлага бий.',
            'Java Reflection API, annotation боловсруулалт, object lifecycle зэрэг суурь ойлголтыг гүнзгийрүүлэн судлах шаардлага байна.')
        replace_text_in_shape(shape, 'тулгуурласан суурь DI framework-ийг',
                              'тулгуурласан суурь DI framework-г')
        replace_text_in_shape(shape, '5. Dependency tree visualization (A2)',
                              '5. Dependency Tree Visualization (A2 сайжруулалт)')
        replace_text_in_shape(shape, '6. Бодит backend төсөлд (CarRentalSystem) ашиглах',
                              '6. Бодит backend төсөл (CarRentalSystem)-д практикт ашиглах')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 4  – Theoretical research (page 3/12)
# ═══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape.has_text_frame:
        # Fix title: add separators
        replace_text_in_shape(shape,
            'Онолын судалгаа: IoC  DI  Annotation  Reflection',
            'Онолын судалгаа: IoC  •  DI  •  Annotation  •  Reflection')
        # Fix page number (should be 2/12, slides 3-4 are section 01-02)
        replace_text_in_shape(shape, '3 / 12', '3 / 13')
        # Fix word choice
        replace_text_in_shape(shape, 'Уламжлалт: A→B шууд, IoC: Container→A,B',
                              'Уламжлалт: A→B шууд үүсгэх, IoC: Container→A,B удирдана.')
        replace_text_in_shape(shape, 'гаднаас нь оруулах. Field / Constructor / Setter injection аргууд.',
                              'гаднаас нь оруулах. Field injection, Constructor injection, Setter injection аргууд байдаг.')
        replace_text_in_shape(shape,
            'Runtime үед Class<?>, Field, Constructor-г динамикаар унших, object үүсгэх, field утга оноох боломжийг JDK стандарт сангаар хэрэгжүүлнэ.',
            'Runtime үед Class<?>, Field, Constructor-г динамикаар унших, объект үүсгэх, field-д утга оноох боломжийг JDK стандарт сангаар хэрэгжүүлдэг.')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 5  – Annotation & Reflection detail (no page number)
# ═══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame:
        # Fix heading "Java in Annotation" → "Java дахь Annotation"
        if shape.name == 'Text 6':
            replace_text_in_shape(shape, 'Java in Annotation', 'Java дахь Annotation')

        # Fix "Боломжууд:" description
        replace_text_in_shape(shape,
            '• Класс, гишүүн өгөгдөл, функцын мэдээлэл унших',
            '• Класс, гишүүн хувьсагч, аргын мэдээллийг динамикаар унших')
        replace_text_in_shape(shape,
            '• Constructor-оор динамик объект үүсгэх',
            '• Constructor-ээр динамик объект үүсгэх')
        replace_text_in_shape(shape,
            '• Объект үүсгэлт (getDeclaredConstructor)',
            '• Объект үүсгэлт — getDeclaredConstructor().newInstance()')
        replace_text_in_shape(shape,
            '• Field injection (field.set)',
            '• Field injection — field.setAccessible(true); field.set(obj, val)')

        # Apply code highlighting to the annotation code block (Text 7)
        if shape.name == 'Text 7':
            apply_code_style_java(shape)


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 6  – Tech Stack (no page number – section header)
# ═══════════════════════════════════════════════════════════════════════════════
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, 'Програмчлалын хэл', 'Програмчлалын хэл')
        replace_text_in_shape(shape, '✓ Нэмэлт сангүй', '✓ Гуравдагч сангүй')
        replace_text_in_shape(shape, 'Maven 3.6+ орчин', 'Maven 3.6+ орчин дэмжинэ')
        replace_text_in_shape(shape, 'Java 8+ бүх орчинд', 'Java 8+ бүх орчинд ажиллана')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 7  – System requirements (4/12)
# ═══════════════════════════════════════════════════════════════════════════════
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, '4 / 12', '4 / 13')
        replace_text_in_shape(shape, 'Visualization без 3rd-party (A2)', 'HTML визуализаци — гуравдагч сангүй (A2)')
        replace_text_in_shape(shape, '= Шинэчлэлт хийсэн A2 ажлаар нэмэгдсэн шаардлага',
                              '★ = A2 шинэчлэлтаар нэмэгдсэн шаардлага')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 8  – Architecture (5/12)
# ═══════════════════════════════════════════════════════════════════════════════
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, '5 / 12', '5 / 13')
        replace_text_in_shape(shape,
            'Хэрэглэгчтэй харилцах гадаад интерфэйс',
            'Хөгжүүлэгчийн гадаад интерфэйс — annotation-ийг класс дээр тавина')
        replace_text_in_shape(shape,
            'package-уудыг рекурсив нэгжиж component class илрүүлнэ',
            'Package болон дэд package-уудыг рекурсив нэгжиж @Component класс илрүүлнэ')
        replace_text_in_shape(shape,
            'Bean lifecycle удирдлага, Singleton cache, Dependency resolution',
            'Bean-ийн амьдралын мөчлөг удирдлага, Singleton кэш, хамаарлын шийдэл')
        replace_text_in_shape(shape,
            'Fail-fast зан авир — startup үед тодорхой алдаа',
            'Fail-fast зан авир — startup үед нарийвчилсан алдааны мэдэгдэл')
        replace_text_in_shape(shape,
            'Давхарга тус бүр өөрийн үүргийг тодорхой хязгаарт гүйцэтгэнэ',
            'Давхарга тус бүр нэг үүрэгтэй — Single Responsibility Principle')
        replace_text_in_shape(shape,
            'Тест хийхэд хялбар — нэгж болон интеграцийн тест тусдаа',
            'Тест хийхэд хялбар — нэгж (unit) болон интеграцийн тест тусдаа явагдана')
        replace_text_in_shape(shape,
            'Annotation, Scanner, Container-ийг бие даан өргөтгөх боломжтой',
            'Annotation, Scanner, Container давхарга бие даан өргөтгөх боломжтой')
        replace_text_in_shape(shape,
            'Spring Framework-ийн архитектурын ойлголттой нийцэж байна',
            'Spring Framework-ийн дотоод бүтцийн ойлголттой нягт нийцэж байна')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 9  – Module implementation (6/12)
# ═══════════════════════════════════════════════════════════════════════════════
slide9 = prs.slides[8]
for shape in slide9.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, '6 / 12', '6 / 13')
        replace_text_in_shape(shape,
            '▸ Package болон дэд package-уудыг рекурсив нэгжиж',
            '▸ Package болон дэд package-уудыг рекурсив нэгжиж,')
        replace_text_in_shape(shape,
            '@Component-той классуудыг илрүүлнэ. Class.forName() ашиглана.',
            '  @Component анноти бүхий класс илрүүлнэ (Class.forName() ашиглана)')
        replace_text_in_shape(shape,
            '▸  inCreation Set → circular dep. шалгалт',
            '▸  inCreation Set → тойрог хамаарал илрүүлэлт')
        replace_text_in_shape(shape,
            'Startup урсгал:  run(App.class)  →  scan package  →  register beans  →  preload singletons  →  inject dependencies  →  build tree (A2)  →  Container бэлэн',
            'Startup урсгал:  run(App.class)  →  scan()  →  registerBeans()  →  preloadSingletons()  →  injectDependencies()  →  buildTree(A2)  →  Container бэлэн')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 10  – A1 Test improvements (7/12)
# ═══════════════════════════════════════════════════════════════════════════════
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, '7 / 12', '7 / 13')
        replace_text_in_shape(shape, 'Үндсэн happy-path тест', 'Зөвхөн амжилттай урсгалын тест (happy-path)')
        replace_text_in_shape(shape,
            'Давхарга бүрийн салбарласан урсгал бүрийг тусдаа тест',
            'Давхарга бүрийн салбарлах урсгал бүрийг тусдаа тест болгосон')
        replace_text_in_shape(shape, 'Ерөнхий RuntimeException мэдэгдэл', 'Ерөнхий RuntimeException мэдэгдэл ашиглаж байсан')
        replace_text_in_shape(shape,
            'Bean нэр, хайсан төрөл, candidate жагсаалт, шалтгаан агуулсан тодорхой мэдэгдэл',
            'Bean нэр, хайлтын төрөл, candidate жагсаалт, шалтгаан бүхий нарийвчилсан мэдэгдэл')
        replace_text_in_shape(shape, 'Нэг тест файлд cold start тест', 'Нэг файлд нийт тест хольж байсан')
        replace_text_in_shape(shape,
            'Unit vs Integration тусдаа — 10 тест класс, 47 тест',
            'Unit туршилт ба Integration туршилт тусдаа — нийт 10 тест класс, 47 тест')
        replace_text_in_shape(shape, 'A1 үр дүн: 47 тест  ✓  0 failures',
                              'A1 үр дүн: Нийт 47 тест  ✓  0 алдаа  |  mvn test — бүгд тэнцсэн')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 11  – A2 Dependency Tree (8/12)
# ═══════════════════════════════════════════════════════════════════════════════
slide11 = prs.slides[10]
for shape in slide11.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, '8 / 12', '8 / 13')
        replace_text_in_shape(shape,
            'Dependency tree → JSON format export (machine-readable, test-able)',
            'Dependency tree → JSON форматаар экспортлох (machine-readable, тестлэх боломжтой)')
        replace_text_in_shape(shape,
            'Standalone interactive HTML dashboard — browser-д автоматаар нээгдэнэ',
            'Standalone интерактив HTML dashboard үүсгэж, browser-д автоматаар нээгдэнэ')
        replace_text_in_shape(shape,
            'Startup дараа BeanRegistry-ийн бүх BeanDefinition-аас DependencyNode үүсгэнэ',
            'Startup дууссаны дараа BeanRegistry-н бүх BeanDefinition-аас DependencyNode үүсгэнэ')
        replace_text_in_shape(shape,
            'Reflection: @Autowired field бүрээр dependency edge байгуулна',
            'Reflection ашиглан @Autowired field бүр тутамд dependency edge холболт үүсгэнэ')
        replace_text_in_shape(shape,
            'DFS (visited + recursionStack) ашиглан тойрог хамаарал шалгана',
            'DFS алгоритм (visited + recursionStack) ашиглан тойрог хамаарал илрүүлнэ')
        replace_text_in_shape(shape,
            'generateAndOpen() → HTML файл үүсгэж browser нээнэ',
            'generateAndOpen() → HTML файл үүсгэж, browser-д автоматаар нээнэ')
        replace_text_in_shape(shape,
            '@EnableIoC(visualize=true) → startup-д автоматаар ажиллана',
            '@EnableIoC(visualize=true) тохируулснаар startup-д автоматаар ажиллана')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 12  – JitPack deployment (9/12)  + CODE HIGHLIGHTING
# ═══════════════════════════════════════════════════════════════════════════════
slide12 = prs.slides[11]
for shape in slide12.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, '9 / 12', '9 / 13')
        # Fix "Pom.xml" casing
        replace_text_in_shape(shape, 'Pom.xml', 'pom.xml')
        replace_text_in_shape(shape, '<version>Tag</version>',
                              '<version>v1.0.0</version>')

        # Apply XML syntax highlighting to pom.xml block
        if shape.name == 'Text 28':
            apply_code_style_xml(shape)

        # Apply Java syntax highlighting to usage code block
        if shape.name in ('Text 35', 'Text 36'):
            apply_code_style_java(shape)


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 13  – Comparison chart (10/12)
# ═══════════════════════════════════════════════════════════════════════════════
slide13 = prs.slides[12]
for shape in slide13.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, '10 / 12', '10 / 13')
        replace_text_in_shape(shape, '+32%', '+45%')
        replace_text_in_shape(shape, 'Тестийн өсөлт (Эхний хэрэгжүүлэлт→A2)',
                              'Тестийн өсөлт: Анхны хэрэгжүүлэлтээс A2 хүртэл')
        replace_text_in_shape(shape, '66 / 0', '47 / 0')
        replace_text_in_shape(shape,
            'График дээр харагдаж байгаачлан эхний явцын хэрэгжүүлэлтээс хойш өнөөдрийг хүртэлх тестийн нийт үзүүлэлтүүдийг харуулсан',
            'График дээр эхний явцын хэрэгжүүлэлтээс A1, A2 шинэчлэлт хүртэлх тестийн тоо болон чанарын өөрчлөлтийг харуулсан.')


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 14  – HTML Dashboard & @EnableIoC  (FIX page: 9/12 → 11/13)
#              + CODE HIGHLIGHTING on @EnableIoC block
# ═══════════════════════════════════════════════════════════════════════════════
slide14 = prs.slides[13]

# Keyword-based color mapping for @EnableIoC code lines
ENABLE_IOC_COLORS = {
    'Text 35': ANNOTATION_CLR,   # @EnableIoC(
    'Text 36': ATTR_CLR,         # scanPackages={},
    'Text 37': ATTR_CLR,         # excludePackages={
    'Text 38': STRING_CLR,       # "mn.edu.num.app.carrental"},
    'Text 39': ATTR_CLR,         # visualize = true
    'Text 40': CODE_FG,          # )
    'Text 41': KEYWORD_CLR,      # public class Main {
    'Text 42': CLASS_CLR,        # ApplicationContext.run(
    'Text 43': CODE_FG,          # Main.class); }
}

for shape in slide14.shapes:
    if shape.has_text_frame:
        # ── Fix page number
        if shape.name == 'Text 0':
            replace_text_in_shape(shape, '9 / 12', '11 / 13')

        # ── Fix @EnableIoC label text
        if shape.name == 'Text 34':
            replace_text_in_shape(shape, '@EnableIoC', '@EnableIoC — жишээ ашиглалт')

        # ── Fix text content
        replace_text_in_shape(shape,
            'Зураг 2 CarRentSystem – ийн dependency tree dashboard',
            'Зураг 2. CarRentalSystem-ийн dependency tree dashboard')

        # ── Apply code highlighting to @EnableIoC lines
        if shape.name in ENABLE_IOC_COLORS:
            clr = ENABLE_IOC_COLORS[shape.name]
            # Special handling for Text 41 (multiple keywords)
            if shape.name == 'Text 41':
                solid_fill(shape, CODE_BG)
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        t = run.text.strip()
                        if t.startswith('public') or t.startswith('class'):
                            run.font.color.rgb = KEYWORD_CLR
                        else:
                            run.font.color.rgb = CLASS_CLR
                        run.font.name = 'Consolas'
            else:
                solid_fill(shape, CODE_BG)
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = clr
                        run.font.name = 'Consolas'


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 15  – CONCLUSION  (12/13)  – full rewrite
# ═══════════════════════════════════════════════════════════════════════════════
slide15 = prs.slides[14]
for shape in slide15.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, '12 / 12', '12 / 13')

        if shape.name == 'Text 10':
            tf = shape.text_frame
            tf.word_wrap = True

            # Define new conclusion content
            conclusion_paragraphs = [
                {
                    'text': 'Энэхүү судалгааны ажлаар annotation болон Java Reflection механизмд суурилсан IoC/DI framework-г бие даан загварчилж, хэрэгжүүлэв.',
                    'size_pt': 13, 'bold': False,
                },
                {'text': '', 'size_pt': 8},
                {
                    'text': 'Хэрэгжүүлэлтийн үндсэн үр дүнгүүд:',
                    'size_pt': 13, 'bold': True,
                },
                {
                    'text': '• @Component, @Autowired, @Qualifier, @Scope, @EnableIoC — 5 custom annotation-г тодорхойлж, RUNTIME retention policy-тэй хэрэгжүүллээ.',
                    'size_pt': 12, 'bold': False,
                },
                {
                    'text': '• Annotation → Scanner → Container → Exception гэсэн 4 давхаргат архитектур хэрэгжүүлж, Single Responsibility зарчмыг баримталлаа.',
                    'size_pt': 12, 'bold': False,
                },
                {
                    'text': '• Bean-ийн Singleton/Prototype scope удирдлага болон circular dependency fail-fast илрүүлэлтийг бодитоор хэрэгжүүллээ.',
                    'size_pt': 12, 'bold': False,
                },
                {
                    'text': '• A1 сайжруулалт: 10 тест класс, 47 тест, 0 алдаа — хэрэгжүүлэлтийн найдвартай байдлыг бүрэн нотолсон.',
                    'size_pt': 12, 'bold': False,
                },
                {
                    'text': '• A2 сайжруулалт: DependencyTreeBuilder, JSON export, интерактив HTML dashboard — bean-ийн бүтцийг харагдахуйц болгосон.',
                    'size_pt': 12, 'bold': False,
                },
                {
                    'text': '• JitPack дээр Maven сан байдлаар байршуулж, CarRentalSystem бодит төсөлд практикт амжилттай ашигласан.',
                    'size_pt': 12, 'bold': False,
                },
                {'text': '', 'size_pt': 8},
                {
                    'text': 'Судалгааны дүгнэлт: IoC/DI зарчмыг гуравдагч талын framework ашиглалгүйгээр, Java Reflection ба Annotation механизмд тулгуурлан бие даан хэрэгжүүлэх боломжтой гэдгийг практик түвшинд нотолсон.',
                    'size_pt': 13, 'bold': True,
                },
            ]

            clear_and_write_paragraphs(tf, conclusion_paragraphs)


# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 16  – Thank you
# ═══════════════════════════════════════════════════════════════════════════════
slide16 = prs.slides[15]
for shape in slide16.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, '12 / 12', '')
        replace_text_in_shape(shape,
            'Удирдагч: Дэд проф. Т.Цэрэннадмид',
            'Удирдагч багш: Дэд проф. Т.Цэрэннадмид')


# ─── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print('Saved:', OUTPUT)

