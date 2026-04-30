from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

INPUT = r"C:\4_Havar\Diplom\Example_project\CarRentalSystem\report\ReportPPT_improved.pptx"
OUTPUT = r"C:\4_Havar\Diplom\Example_project\CarRentalSystem\report\ReportPPT_improved_with_framework_comparison.pptx"

BG = RGBColor(0x08, 0x12, 0x20)
PANEL = RGBColor(0x11, 0x20, 0x37)
HEADER = RGBColor(0x1C, 0x35, 0x58)
ACCENT = RGBColor(0x63, 0xB3, 0xFF)
TEXT = RGBColor(0xE8, 0xF0, 0xFF)
MUTED = RGBColor(0xA8, 0xBA, 0xD6)
BORDER = RGBColor(0x2D, 0x4D, 0x76)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xF8, 0x71, 0x71)

ROWS = [
    (
        "Spring Framework / Spring Boot",
        "Хамгийн өргөн хэрэглээтэй, @Component/@Autowired зэрэг DI концепцийг түгээмэл болгосон.",
        "Heavyweight; startup удаан; тохиргоо ба abstraction их тул сурахад төвөгтэй.",
    ),
    (
        "Google Guice",
        "Хөнгөн, зөвхөн DI-д төвлөрсөн, JSR-330 @Inject ашигладаг.",
        "Автомат package scan бараг байхгүй; bindings-ийг Module дотор гараар өгнө; алдаа runtime дээр илэрч болно.",
    ),
    (
        "Dagger 2",
        "Reflection-гүй compile-time injection хийдэг тул runtime хурдан.",
        "Boilerplate код их; Component/Module тохиргоо төвөгтэй; build хугацаа уртасдаг.",
    ),
    (
        "Jakarta EE CDI / Weld",
        "Java EE/Jakarta EE-ийн албан ёсны DI стандарт, enterprise орчинд нийцтэй.",
        "Application server эсвэл тусгай runtime орчинд илүү тохиромжтой; standalone жижиг төсөлд нүсэр.",
    ),
]


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height, text, *, font_size=18, bold=False,
               color=TEXT, align=PP_ALIGN.LEFT, name=None):
    shape = slide.shapes.add_textbox(left, top, width, height)
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Aptos'
    return shape


def style_rect(shape, fill_rgb, line_rgb=BORDER, radius=True):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(1.2)


def add_panel(slide, left, top, width, height, fill_rgb=PANEL):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    style_rect(shape, fill_rgb)
    return shape


def replace_all_text(prs, old, new):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)


def move_last_slide_to(prs, target_index):
    sld_id_lst = prs.slides._sldIdLst
    last = sld_id_lst[-1]
    sld_id_lst.remove(last)
    sld_id_lst.insert(target_index, last)


def add_framework_slide(prs):
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)

    for shape in list(slide.shapes):
        el = shape.element
        el.getparent().remove(el)

    slide_width = prs.slide_width

    add_textbox(
        slide,
        Inches(0.55), Inches(0.35), Inches(8.5), Inches(0.55),
        "Ижил төрлийн Java DI Framework-үүдийн товч харьцуулалт",
        font_size=24, bold=True,
    )
    add_textbox(
        slide,
        Inches(0.58), Inches(0.9), Inches(11.9), Inches(0.55),
        "Манай custom DI framework нь lightweight, auto-scan, standalone startup гэсэн зорилгоор дээрх framework-үүдийн хүндрэлтэй талыг багасгахыг зорьсон.",
        font_size=11, color=MUTED,
    )

    table_left = Inches(0.55)
    table_top = Inches(1.4)
    table_width = slide_width - Inches(1.1)
    header_height = Inches(0.48)
    row_height = Inches(1.18)

    add_panel(slide, table_left, table_top, table_width, header_height, HEADER)
    framework_x = table_left + Inches(0.16)
    strength_x = table_left + Inches(3.15)
    weakness_x = table_left + Inches(7.4)

    add_textbox(slide, framework_x, table_top + Inches(0.06), Inches(2.7), Inches(0.3), "Framework", font_size=11, bold=True)
    add_textbox(slide, strength_x, table_top + Inches(0.06), Inches(4.0), Inches(0.3), "Товч онцлог", font_size=11, bold=True)
    add_textbox(slide, weakness_x, table_top + Inches(0.06), Inches(4.9), Inches(0.3), "Гол сул тал", font_size=11, bold=True)

    for index, (framework, strength, weakness) in enumerate(ROWS):
        top = table_top + header_height + Inches(0.08) + row_height * index
        panel = add_panel(slide, table_left, top, table_width, row_height, PANEL)
        if index % 2 == 1:
            panel.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2F)

        add_textbox(slide, framework_x, top + Inches(0.1), Inches(2.7), Inches(0.92), framework, font_size=11, bold=True)
        add_textbox(slide, strength_x, top + Inches(0.1), Inches(4.0), Inches(0.92), strength, font_size=10.5, color=TEXT)
        add_textbox(slide, weakness_x, top + Inches(0.1), Inches(4.9), Inches(0.92), weakness, font_size=10.5, color=TEXT)

    summary = add_panel(slide, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.62), RGBColor(0x0E, 0x27, 0x22))
    summary.line.color.rgb = RGBColor(0x1D, 0x7A, 0x63)
    add_textbox(
        slide,
        Inches(0.9), Inches(6.58), Inches(11.6), Inches(0.28),
        "Дүгнэлт: Custom framework нь сурахад хялбар, dependency wiring-ийг ойлгоход илүү тунгалаг, жижиг/дунд хэмжээний standalone backend төсөлд тохиромжтой.",
        font_size=11, color=RGBColor(0xCC, 0xFB, 0xF1), bold=True,
    )

    add_textbox(slide, Inches(11.6), Inches(7.0), Inches(1.0), Inches(0.28), "12 / 14", font_size=11, color=MUTED, align=PP_ALIGN.RIGHT)

    return slide


def main():
    prs = Presentation(INPUT)

    replace_all_text(prs, ' / 13', ' / 14')
    replace_all_text(prs, '12 / 14', '13 / 14')

    add_framework_slide(prs)
    move_last_slide_to(prs, 14)

    prs.save(OUTPUT)
    print(f'Saved: {OUTPUT}')
    print(f'Slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()

