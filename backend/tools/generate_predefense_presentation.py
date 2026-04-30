from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

ROOT = Path(r"C:\4_Havar\Diplom\Example_project\CarRentalSystem")
REPORT_DIR = ROOT / "report"
OUTPUT = REPORT_DIR / "DI_Framework_PreDefense_CarRental_2026.pptx"

BG = RGBColor(7, 17, 34)
PANEL = RGBColor(15, 28, 51)
PANEL_ALT = RGBColor(20, 40, 74)
ACCENT = RGBColor(53, 129, 255)
ACCENT_2 = RGBColor(19, 194, 194)
TEXT = RGBColor(239, 244, 255)
MUTED = RGBColor(178, 192, 214)
SUCCESS = RGBColor(59, 201, 123)
WARNING = RGBColor(255, 196, 61)
DANGER = RGBColor(255, 122, 122)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.1), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Arial"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = TEXT
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.03), Inches(11.8), Inches(0.4))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Arial"
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = MUTED


def add_footer(slide, idx, total):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.52), Inches(6.92), Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = PANEL_ALT
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(11.85), Inches(6.98), Inches(0.75), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{idx}/{total}"
    r.font.name = "Arial"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_panel(slide, left, top, width, height, title=None, fill=PANEL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = PANEL_ALT
    shape.line.width = Pt(1)
    if title:
        tx = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.3))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Arial"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE
    return shape


def add_bullets(slide, left, top, width, height, bullets, font_size=17, color=TEXT, level_step=0.3):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, level, bullet_color = item
        else:
            text, level, bullet_color = item, 0, color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.font.name = "Arial"
        p.font.size = Pt(font_size - level * level_step * 10)
        p.font.color.rgb = bullet_color
        p.line_spacing = 1.15
        p.space_after = Pt(5)
        p.bullet = True
    return box


def add_label_value(slide, left, top, width, label, value, accent=ACCENT):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.82))
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = accent
    box.line.width = Pt(1)
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.65))
    tf = tb.text_frame
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = label
    r1.font.name = "Arial"
    r1.font.size = Pt(10)
    r1.font.color.rgb = MUTED
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    r2.font.name = "Arial"
    r2.font.size = Pt(16)
    r2.font.bold = True
    r2.font.color.rgb = WHITE


def add_table(slide, left, top, width, height, columns, rows):
    table = slide.shapes.add_table(len(rows) + 1, len(columns), left, top, width, height).table
    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL_ALT
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Arial"
                r.font.size = Pt(11)
                r.font.bold = True
                r.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "Arial"
                    r.font.size = Pt(10.5)
                    r.font.color.rgb = TEXT
    return table


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []

    # 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.48), Inches(12.45), Inches(5.95))
    hero.fill.solid()
    hero.fill.fore_color.rgb = PANEL
    hero.line.color.rgb = PANEL_ALT
    hero.line.width = Pt(1.2)
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(0.95), Inches(2.7), Inches(0.45))
    chip.fill.solid(); chip.fill.fore_color.rgb = ACCENT; chip.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(1.05), Inches(1.02), Inches(2.3), Inches(0.25))
    p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = "Урьдчилсан хамгаалалтын илтгэл"; r.font.name = "Arial"; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
    title = slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(8.7), Inches(1.8))
    tf = title.text_frame
    for idx, line in enumerate([
        "Backend системд зориулсан",
        "annotation-д суурилсан Java framework-ийн",
        "загварчлал, хэрэгжүүлэлт"
    ]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line; r.font.name = "Arial"; r.font.size = Pt(24 if idx == 0 else 22); r.font.bold = True; r.font.color.rgb = TEXT
        p.space_after = Pt(2)
    sub = slide.shapes.add_textbox(Inches(0.92), Inches(3.6), Inches(8.6), Inches(1.0))
    tf = sub.text_frame
    p = tf.paragraphs[0]; r = p.add_run(); r.text = "DI framework-ийг бодит Car Rental System төсөлд ашиглаж, backend + frontend + PostgreSQL түвшинд баталгаажуулсан шинэчилсэн хувилбар"; r.font.name = "Arial"; r.font.size = Pt(15); r.font.color.rgb = MUTED
    add_label_value(slide, Inches(0.9), Inches(5.05), Inches(3.0), "Гүйцэтгэсэн", "Г.Батням / 22B1NUM5578")
    add_label_value(slide, Inches(4.05), Inches(5.05), Inches(2.95), "Удирдагч", "Др. Т.Цэрэннадмид", ACCENT_2)
    add_label_value(slide, Inches(7.18), Inches(5.05), Inches(2.1), "Огноо", "2026.04", SUCCESS)
    add_label_value(slide, Inches(9.45), Inches(5.05), Inches(2.95), "Судалгааны шат", "12-р урьдчилсан хамгаалалт", WARNING)
    note = slide.shapes.add_textbox(Inches(9.6), Inches(1.5), Inches(2.4), Inches(2.6))
    tf = note.text_frame
    for idx, t in enumerate(["Шинэ төвлөрөл", "• Framework + бодит төсөл", "• CRUD UI ба тусдаа хуудаснууд", "• PostgreSQL schema init", "• 5 backend test + frontend build"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph(); r = p.add_run(); r.text = t; r.font.name = "Arial"; r.font.size = Pt(15 if idx == 0 else 13); r.font.bold = idx == 0; r.font.color.rgb = WHITE if idx == 0 else MUTED
    slides.append(slide)

    # 2
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "Судалгааны үндэслэл, зорилго, шинэчлэл", "Тайлангийн удиртгал + хэрэглээний ажлын өнөөгийн хүрээ")
    add_panel(slide, Inches(0.55), Inches(1.45), Inches(4.0), Inches(4.95), "Асуудал")
    add_bullets(slide, Inches(0.78), Inches(1.9), Inches(3.5), Inches(4.2), [
        "Framework-ийг ашигладаг ч дотоод ажиллах зарчмыг практикт бүрэн ойлгохгүй байх асуудал түгээмэл.",
        "IoC, DI, annotation, reflection зэрэг ойлголтуудыг дан онолоор тайлбарлах нь хангалтгүй.",
        "Бодит backend төсөл дээр ажиллуулж байж container-ийн үнэ цэнэ, хязгаар тод харагдана."
    ], font_size=15)
    add_panel(slide, Inches(4.7), Inches(1.45), Inches(4.0), Inches(4.95), "Зорилго")
    add_bullets(slide, Inches(4.93), Inches(1.9), Inches(3.52), Inches(4.2), [
        "Annotation-д суурилсан DI framework-ийн цөм ажиллагааг хэрэгжүүлэх.",
        "Framework-ийг CarRentalSystem backend-д нэгтгэн бодит хэрэглээнд турших.",
        "Вэб API, өгөгдлийн сан, UI урсгалтай хамтатган эцсийн үр дүнг баталгаажуулах."
    ], font_size=15)
    add_panel(slide, Inches(8.85), Inches(1.45), Inches(3.95), Inches(4.95), "Энэ удаагийн шинэ фокус", fill=PANEL_ALT)
    add_bullets(slide, Inches(9.08), Inches(1.9), Inches(3.45), Inches(4.2), [
        "Өмнөх илтгэлд тайлбарласан онолыг бодит төсөлд хэрэглэсэн.",
        "CRUD backend API, React frontend, PostgreSQL persistence нэмсэн.",
        "UI/UX сайжруулалт, тусдаа edit/rent page, build/test үр дүнг илтгэлд тусгав."
    ], font_size=15)
    slides.append(slide)

    # 3
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "Тайлангаас авсан шаардлага ба зохиомж", "main.pdf-ийн 2, 3-р бүлгийн гол санааг товчлон нэгтгэв")
    add_panel(slide, Inches(0.55), Inches(1.45), Inches(4.0), Inches(5.1), "Функциональ шаардлага")
    add_bullets(slide, Inches(0.78), Inches(1.92), Inches(3.5), Inches(4.4), [
        "Package scanning ба @Component илрүүлэлт",
        "Bean registration, object instantiation, getBean() retrieval",
        "@Autowired injection, @Qualifier сонголт, scope удирдлага",
        "Fail-fast exception handling, dependency tree visualization"
    ], font_size=15)
    add_panel(slide, Inches(4.7), Inches(1.45), Inches(4.0), Inches(5.1), "Функциональ бус шаардлага")
    add_bullets(slide, Inches(4.93), Inches(1.92), Inches(3.5), Inches(4.4), [
        "Java SE-д тулгуурласан, энгийн API-тай байх",
        "Startup үеийн найдвартай ажиллагаа, circular dependency илрүүлэх",
        "Тестлэгдэх, өргөтгөх, ойлгомжтой лог ба бүтцээр хангах",
        "Гуравдагч этгээдийн хүнд framework-ээс хамаарал багатай байх"
    ], font_size=15)
    add_panel(slide, Inches(8.85), Inches(1.45), Inches(3.95), Inches(5.1), "Сонгосон зохиомж")
    add_bullets(slide, Inches(9.08), Inches(1.92), Inches(3.45), Inches(4.4), [
        "Давхаргат / layered architecture",
        "annotation → scanner → container → exception багцын хуваарилалт",
        "ApplicationContext-ийг төв удирдлага болгон ашиглах",
        "Тайлангийн package diagram, activity/sequence зураглалтай уялдуулсан"
    ], font_size=15)
    slides.append(slide)

    # 4
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "DI framework-ийн үндсэн архитектур", "Өмнөх presentation-ийн онолын хэсгийг хэрэгжилттэй холбон шинэчилсэн")
    add_panel(slide, Inches(0.55), Inches(1.35), Inches(12.2), Inches(1.2), "Startup flow")
    flow = [
        ("1. scan", Inches(0.95), ACCENT),
        ("2. register", Inches(3.35), ACCENT_2),
        ("3. singleton preload", Inches(5.78), WARNING),
        ("4. inject", Inches(8.55), SUCCESS),
        ("5. ready context", Inches(10.75), DANGER),
    ]
    for text, x, color in flow:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(1.75), Inches(1.75), Inches(0.48))
        box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
        t = slide.shapes.add_textbox(x + Inches(0.08), Inches(1.84), Inches(1.58), Inches(0.2))
        p = t.text_frame.paragraphs[0]; r = p.add_run(); r.text = text; r.font.name = "Arial"; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE
    add_panel(slide, Inches(0.55), Inches(2.85), Inches(5.9), Inches(3.55), "Framework core")
    add_bullets(slide, Inches(0.82), Inches(3.3), Inches(5.35), Inches(2.8), [
        "@EnableIoC + ApplicationContext.run(App.class) хэлбэрээр startup туршлага бий болгосон.",
        "@Component, @Autowired, @Qualifier, @Scope annotation-ууд runtime үед reflection-оор уншигдана.",
        "Dependency tree HTML visualization нь startup lifecycle-тэй интеграцлагдсан.",
        "Fail-fast exception-ууд: bean creation, missing bean, circular dependency нөхцөлийг оношлоно."
    ], font_size=15)
    add_panel(slide, Inches(6.65), Inches(2.85), Inches(6.1), Inches(3.55), "Current evidence from project")
    add_bullets(slide, Inches(6.92), Inches(3.3), Inches(5.55), Inches(2.8), [
        "IoC container нь 16+ component-ийг илрүүлж inject хийж байна.",
        "dependency-tree.html файл backend startup/test үед автоматаар үүсэж байна.",
        "Servlet, service, persistence, config зэрэг давхаргууд framework-ийн bean байдлаар ажиллаж байна.",
        "Framework нь зөвхөн demo биш, backend application wiring-д бодитоор ашиглагдав."
    ], font_size=15)
    slides.append(slide)

    # 5
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "CarRentalSystem-д хийсэн бодит интеграц", "Framework-ийг ашигласан backend системийн бүтэц")
    cols = ["UI", "Frontend", "Backend API", "Core Service", "Persistence"]
    xs = [0.75, 3.05, 5.45, 7.95, 10.4]
    colors = [ACCENT, ACCENT_2, WARNING, SUCCESS, DANGER]
    labels = [
        "Хэрэглэгч\nDashboard / CRUD",
        "React + Vite\nRouter + Toast + Theme",
        "Embedded Tomcat\nCarServlet / RentalServlet",
        "CarRentalService\nUse case orchestration",
        "PostgreSQL + fallback\nSchema init"
    ]
    for x, c, label in zip(xs, colors, labels):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.0), Inches(2.0))
        shape.fill.solid(); shape.fill.fore_color.rgb = c; shape.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(x)+Inches(0.12), Inches(2.23), Inches(1.76), Inches(1.5))
        tf = tb.text_frame
        for idx, line in enumerate(label.split("\n")):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = line; r.font.name = "Arial"; r.font.size = Pt(16 if idx == 0 else 13); r.font.bold = True if idx == 0 else False; r.font.color.rgb = WHITE
    add_panel(slide, Inches(0.72), Inches(4.55), Inches(12.0), Inches(1.65), "Интеграцын үр дүн")
    add_bullets(slide, Inches(0.95), Inches(4.95), Inches(11.45), Inches(0.95), [
        "Domain: Car, Rental, Invoice • Ports: in/out use cases • Adapters: web + persistence • Config: ApplicationConfig, DatabaseConfig",
        "IRentCarUseCase, IManageCarsUseCase, IGetCarsUseCase зэрэг интерфэйсээр service layer тусгаарлагдсан.",
        "Энэ бүтэц нь тайлангийн layered architecture сонголтыг бодит төслөөр баталж байна."
    ], font_size=14)
    slides.append(slide)

    # 6
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "Backend хэрэгжүүлэлт: API, бизнес логик, persistence", "Одоогийн хамгаалалтын гол нэмэлт нь бодит системийн ажиллагаа")
    add_panel(slide, Inches(0.55), Inches(1.45), Inches(4.05), Inches(4.95), "REST API")
    add_bullets(slide, Inches(0.82), Inches(1.92), Inches(3.55), Inches(4.2), [
        "GET /api/cars, GET /api/cars/{id}, GET /api/cars/available",
        "POST /api/cars, PUT /api/cars/{id}, DELETE /api/cars/{id}",
        "GET /api/rentals, POST /api/rentals, GET /health",
        "Tomcat embedded server + CORS filter + JSON serialization"
    ], font_size=14.5)
    add_panel(slide, Inches(4.75), Inches(1.45), Inches(4.05), Inches(4.95), "Service дүрэм")
    add_bullets(slide, Inches(5.02), Inches(1.92), Inches(3.55), Inches(4.2), [
        "Машин нэмэх, шинэчлэх, устгах, түрээслэх урсгал",
        "Plate number давхардал шалгах, invalid rental type хамгаалах",
        "Rental history байгаа машины устгалтыг block хийх",
        "Invoice үүсгэх, availability төлөв шинэчлэх"
    ], font_size=14.5)
    add_panel(slide, Inches(8.95), Inches(1.45), Inches(3.85), Inches(4.95), "Persistence")
    add_bullets(slide, Inches(9.2), Inches(1.92), Inches(3.35), Inches(4.2), [
        "PostgreSQL adapter + fallback in-memory store",
        "schema.sql-аар cars, rentals хүснэгтүүдийг init хийх",
        "Environment override: DB_URL, DB_USERNAME, DB_PASSWORD",
        "PostgreSQL driver 42.7.4 ашигласан"
    ], font_size=14.5)
    slides.append(slide)

    # 7
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "Frontend шинэчлэл ба CRUD урсгал", "UI-г хэрэглэгчийн ажиллагаанд төвлөрүүлж, тусдаа хуудсуудтай болгосон")
    add_panel(slide, Inches(0.55), Inches(1.45), Inches(5.95), Inches(4.9), "Шинэ UI бүтэц")
    add_bullets(slide, Inches(0.82), Inches(1.92), Inches(5.4), Inches(4.2), [
        "React Router ашиглан Dashboard, Cars, CarEdit, CarRent гэсэн page бүтэц гаргасан.",
        "Cars page-д жагсаалт, хайлт, delete; edit болон rent үйлдлийг тусдаа route-р нээдэг болгосон.",
        "Toast notification, dark/light theme toggle, summary cards, dashboard quick action нэмсэн.",
        "Хэт нуршсан backend status текстүүдийг цэвэрлэж, UX-ийг илүү товч, хамгаалалтанд танилцуулахад тохиромжтой болгосон."
    ], font_size=14.5)
    add_panel(slide, Inches(6.7), Inches(1.45), Inches(6.1), Inches(4.9), "CRUD user flow")
    add_bullets(slide, Inches(6.97), Inches(1.92), Inches(5.55), Inches(4.2), [
        "Create: шинэ машин нэмэх form",
        "Read: хайлттай машины жагсаалт",
        "Update: /cars/:carId/edit тусдаа хуудас",
        "Delete: confirm + conflict message",
        "Rent: /cars/:carId/rent тусдаа хуудас, invoice мессеж",
        "Сүүлд layout overflow асуудлуудыг засаж card/header бүрэн харагддаг болгосон"
    ], font_size=14.5)
    slides.append(slide)

    # 8
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "PostgreSQL холболт, schema init, operational readiness", "Өгөгдөл хадгалалт ба runtime бэлэн байдлын тайлбар")
    add_panel(slide, Inches(0.55), Inches(1.45), Inches(4.0), Inches(5.0), "DB integration")
    add_bullets(slide, Inches(0.82), Inches(1.92), Inches(3.5), Inches(4.2), [
        "application.properties дээр db.enabled=true, jdbc:postgresql://localhost:5432/carrental тохиргоо байгаа.",
        "DatabaseInitializer нь startup үед schema.sql-г уншиж хүснэгт үүсгэнэ.",
        "cars seed өгөгдөл ON CONFLICT DO NOTHING байдлаар анхдагчаар орно."
    ], font_size=14.5)
    add_panel(slide, Inches(4.75), Inches(1.45), Inches(4.05), Inches(5.0), "Найдвартай ажиллагаа")
    add_bullets(slide, Inches(5.02), Inches(1.92), Inches(3.55), Inches(4.2), [
        "PostgreSQL унтарсан эсвэл credential буруу үед хэрэглэгчид алдаа тайлбарлагдана.",
        "Repository adapter-ууд fallback store ашиглаж demo/test урсгалыг үргэлжлүүлэх боломжтой.",
        "Test profile дээр db.enabled=false тул unit/integration тест DB-ээс үл хамааран ногоон гарч байна."
    ], font_size=14.5)
    add_panel(slide, Inches(8.95), Inches(1.45), Inches(3.85), Inches(5.0), "Хамгаалалтын үеийн тайлбар")
    add_bullets(slide, Inches(9.2), Inches(1.92), Inches(3.35), Inches(4.2), [
        "pgAdmin дээр database, user, password тохируулж бодит persistence-ийг үзүүлж болно.",
        "DB credential зөв үед rental history restart дараа хадгалагдана.",
        "DB credential буруу үед crash бус, ажиллагааны төлөвийг ойлгомжтой тайлбарлах боломжтой."
    ], font_size=14.5)
    slides.append(slide)

    # 9
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "Туршилт ба баталгаажуулсан үр дүн", "Энэ илтгэлд шууд ашиглах бодит шалгалтын нотолгоо")
    add_label_value(slide, Inches(0.75), Inches(1.45), Inches(2.1), "Backend tests", "5 / 5 pass", SUCCESS)
    add_label_value(slide, Inches(3.0), Inches(1.45), Inches(2.65), "Surefire хугацаа", "7.796 sec", ACCENT)
    add_label_value(slide, Inches(5.8), Inches(1.45), Inches(2.9), "Frontend build", "vite build ✓", ACCENT_2)
    add_label_value(slide, Inches(8.85), Inches(1.45), Inches(3.1), "Bundle", "205.54 kB JS", WARNING)
    add_panel(slide, Inches(0.55), Inches(2.55), Inches(5.95), Inches(3.75), "Backend дээр шалгасан")
    add_bullets(slide, Inches(0.82), Inches(3.02), Inches(5.4), Inches(3.0), [
        "Seeded available cars уншигдаж буйг баталсан.",
        "Rent хийхэд invoice үүсэж, availability жагсаалтаас машин хасагддагийг шалгасан.",
        "Incompatible rental type үед exception шидэгддэг.",
        "Car CRUD use case-ууд амжилттай ажилласан.",
        "Rental history-тэй машиныг delete хийхийг block хийсэн."
    ], font_size=14.5)
    add_panel(slide, Inches(6.7), Inches(2.55), Inches(6.1), Inches(3.75), "Frontend дээр шалгасан")
    add_bullets(slide, Inches(6.97), Inches(3.02), Inches(5.55), Inches(3.0), [
        "Vite production build амжилттай гарсан.",
        "Multi-page routing, CRUD form, search, notification урсгал compile-д асуудалгүй.",
        "UI overflow / half-visible component асуудлуудыг сүүлд зассан.",
        "Ингэснээр хамгаалалтын demo хийхэд frontend тогтвортой болсон."
    ], font_size=14.5)
    slides.append(slide)

    # 10
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "Өмнөх presentation-тай харьцуулсан шинэчлэл", "DI_Framework_Presentation_Updated-д байсан санааг одоогийн биелэлттэй харьцуулсан")
    add_table(slide, Inches(0.55), Inches(1.55), Inches(12.1), Inches(4.95), ["Хэсэг", "Өмнөх presentation", "Одоогийн шинэчилсэн төлөв"], [
        ["Анхаарал төвлөрөл", "DI framework-ийн онол, үндсэн хэрэгжүүлэлт", "DI framework + бодит CarRentalSystem интеграц"],
        ["Архитектур", "Framework-ийн 4 давхарга, custom annotation", "Framework дээр тулгуурласан layered backend + web + persistence"],
        ["Туршилт", "Framework түвшний тест ба coverage", "5 backend integration/use-case test + frontend build + runtime demo бэлэн"],
        ["Цаашдын ажил", "Бодит төсөлд турших, web/data чиглэлээр өргөтгөх", "Бодит төсөлд туршсан, CRUD UI, PostgreSQL, API endpoint-ууд хэрэгжсэн"],
        ["Хэрэглэгчийн түвшин", "Framework developer-centric тайлбар", "Хэрэглэгчийн UI, захиалга, машины удирдлагаар үр дүн харагдана"],
    ])
    note = slide.shapes.add_textbox(Inches(0.75), Inches(6.72), Inches(11.6), Inches(0.35))
    p = note.text_frame.paragraphs[0]; r = p.add_run(); r.text = "Өөрөөр хэлбэл өмнөх илтгэлд 'цаашид бодит төсөлд туршина' гэж байсан санаа энэ удаад хэрэгжиж, хэмжигдсэн үр дүнтэй болсон."; r.font.name = "Arial"; r.font.size = Pt(12.5); r.font.color.rgb = MUTED
    slides.append(slide)

    # 11
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "Үечилсэн төлөвлөгөөтэй уялдуулсан тайлбар", "Үечилсэн_төлөвлөгөө_DI_Framework (1)-ийн 15 долоо хоногийн төлөвлөлттэй холбов")
    add_table(slide, Inches(0.55), Inches(1.55), Inches(12.1), Inches(4.6), ["Төлөвлөгөөний багц ажил", "Тайлан/код дээрх биелэлт", "Одоогийн төлөв"], [
        ["1. Сэдвийн судалгаа", "IoC, DI, annotation, reflection, Spring Core судалгаа main.pdf 1-р бүлэгт бүрэн туссан", "Биелсэн"],
        ["2. Шаардлага, зохиомж", "Framework requirements, package/layer design, diagram-ууд main.pdf 2-3-р бүлэгт туссан", "Биелсэн"],
        ["3. Framework хэрэгжүүлэлт", "Annotation, scanner, container, injector, dependency tree subsystem хэрэгжсэн", "Биелсэн"],
        ["4. Харьцуулалт, туршилт", "Existing framework-үүдтэй харьцуулж, JUnit тестүүд хийсэн", "Биелсэн"],
        ["5. Сайжруулалт, бодит төсөл", "CarRentalSystem backend/frontend, CRUD, PostgreSQL интеграц хийсэн", "Энэ хамгаалалтад гол шинэ үр дүн"],
        ["6. Тайлан боловсруулалт", "Тайлан + шинэчилсэн илтгэлийг урьдчилсан хамгаалалтын шатанд нэгтгэв", "Явагдаж буй / бэлэн"],
    ])
    foot = slide.shapes.add_textbox(Inches(0.75), Inches(6.35), Inches(11.7), Inches(0.5))
    tf = foot.text_frame
    p = tf.paragraphs[0]; r = p.add_run(); r.text = "13-р долоо хоногийн урьдчилсан хамгаалалтын шаардлагад нийцүүлэн: судалгаа → зохиомж → хэрэгжүүлэлт → туршилт → бодит төслийн баталгаажуулалт гэсэн дарааллыг илтгэлд тодорхой харуулав."; r.font.name = "Arial"; r.font.size = Pt(12.5); r.font.color.rgb = MUTED
    slides.append(slide)

    # 12
    slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "Дүгнэлт", "Framework-ийн судалгааг бодит системийн хэрэгжилтээр баталгаажуулсан үр дүн")
    add_panel(slide, Inches(0.78), Inches(1.55), Inches(12.0), Inches(4.05), "Гол дүгнэлт")
    add_bullets(slide, Inches(1.08), Inches(2.0), Inches(11.35), Inches(3.3), [
        "Annotation-д суурилсан DI framework нь backend системийн wiring, component lifecycle, dependency management-ийг ойлгомжтой хэрэгжүүлж чадсаныг бодит төслөөр баталлаа.",
        "CarRentalSystem дээр framework-ийг ашигласнаар layered architecture, CRUD API, embedded Tomcat, PostgreSQL persistence, React frontend зэрэг хэсгүүд нэг шийдэл болж ажиллаж байна.",
        "Туршилтын нотолгоо: backend 5/5 test pass, frontend production build амжилттай, dependency tree visualization автоматаар үүсэж байна.",
        "Хамгаалалтын гол санаа: өмнөх илтгэлд тайлбарласан онолын framework одоо хэрэглээний систем дээр ажиллаж, сайжруулалтын үр дүн нь хэмжигдэж буй шатанд хүрсэн."
    ], font_size=15)
    add_panel(slide, Inches(0.78), Inches(5.95), Inches(5.7), Inches(0.78), "Дараагийн шат")
    add_bullets(slide, Inches(1.02), Inches(6.15), Inches(5.2), Inches(0.4), ["Final defense-д demo flow, pgAdmin persistence proof, боломжтой бол coverage/chart нэмэх"], font_size=13.5)
    add_panel(slide, Inches(6.68), Inches(5.95), Inches(6.1), Inches(0.78), "Баярлалаа")
    add_bullets(slide, Inches(6.92), Inches(6.15), Inches(5.6), Inches(0.4), ["Асуулт, санал хүсэлт авахад бэлэн."], font_size=13.5)
    slides.append(slide)

    total = len(slides)
    for idx, slide in enumerate(slides, start=1):
        add_footer(slide, idx, total)

    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    build_presentation()

